"""Generate the XIII SARL pitch deck (16:9 PPTX).

Run with: python3 generate_xiii_presentation.py
Requires: pip install python-pptx
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'XIII_SARL_PitchDeck.pptx')

ppt = Presentation()
ppt.slide_height = Inches(7.5)
ppt.slide_width = Inches(13.333)  # 16:9

# colors
PANTONE_RGB = (0, 51, 160)  # approx Pantone 286C #0033A0
WHITE = (255, 255, 255)
GRAY = (191, 191, 191)
BLACK = (0, 0, 0)


def add_cover(title, subtitle_lines):
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*BLACK)

    # Title
    tx = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.name = 'Playfair Display'
    p.font.color.rgb = RGBColor(*PANTONE_RGB)

    # Subtitle (one paragraph per line)
    sx = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(1.5))
    sf = sx.text_frame
    sf.word_wrap = True
    for i, line in enumerate(subtitle_lines):
        s = sf.paragraphs[0] if i == 0 else sf.add_paragraph()
        s.text = line
        s.font.size = Pt(20)
        s.font.name = 'Montserrat'
        s.font.color.rgb = RGBColor(*GRAY)


def add_slide(title, bullets):
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*BLACK)

    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.9), Inches(1.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = 'Playfair Display'
    p.font.color.rgb = RGBColor(*PANTONE_RGB)

    # bullets
    bx = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.1), Inches(4.5))
    bf = bx.text_frame
    bf.word_wrap = True
    for i, b in enumerate(bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = '• ' + b
        p.font.size = Pt(18)
        p.font.name = 'Montserrat'
        p.font.color.rgb = RGBColor(*WHITE)
        p.space_after = Pt(12)


add_cover('XIII SARL', ['Culture · Food · Immobilier', 'Expériences curatoriales et culinaires.'])

slides = [
    ("One-liner & Mission", [
        "Créons des expériences exclusives (dîners, ventes d'art, conseil).",
        "Fusionner art et gastronomie pour valeur culturelle et financière.",
    ]),
    ("Problème", [
        "Marché haut de gamme atomisé.",
        "Revenus événements volatils ; besoin d'assise patrimoniale.",
    ]),
    ("Solution", [
        "Dîners éphémères curatorisés + vente d'art.",
        "Conseil stratégique food & branding.",
        "Actif immobilier pour stabiliser cashflow.",
    ]),
    ("Offre / Services", [
        "Dîners, vente d'art, consulting B2B, location immobilière.",
    ]),
    ("Marché & cible", [
        "Clientèle premium LU/FR, entreprises, collectionneurs.",
    ]),
    ("Business model", [
        "Tickets, commissions art, fees consulting, loyers.",
    ]),
    ("Structure juridique", [
        "Holding LU / SARL Opérations / SARL Immobilier / Future LTD HK.",
    ]),
    ("Équipe", [
        "Founder + réseau prestataires (chefs, artistes).",
        "Modèle lean via prestations.",
    ]),
    ("Finances clés", [
        "Apport 250k€, prix bien 650-900k€, emprunt 400-650k€.",
        "Salaires Y1 5k€/mois → Y3 10k€/mois.",
    ]),
    ("Ask", [
        "Finaliser prêt + fonds de roulement 50-150k€.",
        "Partenariats chefs/galeries, budget marketing.",
    ]),
    ("Roadmap", [
        "0-6m : achat + saison 1 dîners.",
        "6-24m : scaler events & consulting.",
        "3-7y : HK expansion.",
    ]),
]

for t, b in slides:
    add_slide(t, b)

ppt.save(OUTPUT_PATH)
print(f"PPTX généré: {OUTPUT_PATH}")
